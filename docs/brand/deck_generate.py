# -*- coding: utf-8 -*-
"""Blades of Anatolia: Ertugrul — loyiha taqdimoti."""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

HERE  = os.path.dirname(os.path.abspath(__file__))
BRAND = r"D:\My_apps\Ertugrul\docs\brand"
SHOTS = os.path.join(os.path.dirname(HERE))

# --- "Temir va Firuza" dizayn tizimi ---
FON    = RGBColor(0x0E, 0x13, 0x16)
PANEL  = RGBColor(0x16, 0x1D, 0x21)
PANEL2 = RGBColor(0x1B, 0x24, 0x28)
CHIZIQ = RGBColor(0x2A, 0x35, 0x3A)
MATN   = RGBColor(0xE4, 0xEA, 0xEA)
SONIK  = RGBColor(0x7C, 0x8B, 0x8F)
FERUZA = RGBColor(0x48, 0xA9, 0xB5)
ZARHAL = RGBColor(0xC0, 0x96, 0x60)
YARA   = RGBColor(0xBC, 0x5A, 0x44)
SUYAK  = RGBColor(0xDC, 0xD3, 0xC4)

H_FONT = "Cambria"      # sarlavha (xavfsiz ro'yxat)
B_FONT = "Calibri"      # matn

W, H = 13.333, 7.5
prs = Presentation()
prs.slide_width  = Inches(W)
prs.slide_height = Inches(H)
BLANK = prs.slide_layouts[6]


def shot(name, folder="final"):
    p = os.path.join(SHOTS, folder, name)
    return p if os.path.exists(p) else None


def new(dark=True):
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = FON if dark else PANEL
    return s


def box(s, x, y, w, h, fill=None, line=None, lw=1.0, shape=MSO_SHAPE.RECTANGLE):
    sp = s.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(lw)
    sp.shadow.inherit = False
    if sp.has_text_frame:
        sp.text_frame.text = ""
    return sp


def mark(s, x, y, size=0.085, col=FERUZA):
    """Dizayn tizimining 'mix boshi' kvadrat motivi."""
    return box(s, x, y, size, size, fill=col)


def text(s, x, y, w, h, runs, size=15, color=MATN, font=B_FONT, bold=False,
         align=PP_ALIGN.LEFT, space=6, line=None, anchor=MSO_ANCHOR.TOP):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    items = runs if isinstance(runs, list) else [runs]
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space)
        if line:
            p.line_spacing = line
        parts = it if isinstance(it, list) else [(it, {})]
        for txt, o in parts:
            r = p.add_run(); r.text = txt
            f = r.font
            f.name = o.get("font", font)
            f.size = Pt(o.get("size", size))
            f.bold = o.get("bold", bold)
            f.italic = o.get("italic", False)
            f.color.rgb = o.get("color", color)
    return tb


def title(s, t, sub=None):
    mark(s, 0.72, 0.66, 0.10, FERUZA)
    text(s, 0.95, 0.52, 11.9, 0.7, t, size=32, color=SUYAK, font=H_FONT, bold=True)
    if sub:
        text(s, 0.95, 1.16, 11.9, 0.42, sub, size=14, color=SONIK)
    return 1.75 if sub else 1.48


def bullets(s, x, y, w, items, size=14, gap=0.40, col=MATN, lead=ZARHAL):
    """Har bir element: (qalin bosh, izoh) yoki oddiy matn."""
    yy = y
    for it in items:
        if isinstance(it, tuple):
            head, body = it
            mark(s, x, yy + 0.09, 0.075, FERUZA)
            text(s, x + 0.22, yy, w - 0.22, gap,
                 [[(head + "  ", {"bold": True, "color": lead}), (body, {"color": col})]],
                 size=size, line=1.15)
        else:
            mark(s, x, yy + 0.09, 0.075, FERUZA)
            text(s, x + 0.22, yy, w - 0.22, gap, it, size=size, color=col, line=1.15)
        yy += gap
    return yy


def card(s, x, y, w, h, head, body, col=ZARHAL, hsize=15, bsize=12):
    box(s, x, y, w, h, fill=PANEL, line=CHIZIQ, lw=0.75)
    mark(s, x + 0.22, y + 0.26, 0.075, FERUZA)
    text(s, x + 0.42, y + 0.18, w - 0.62, 0.32, head, size=hsize, color=col,
         font=H_FONT, bold=True)
    text(s, x + 0.22, y + 0.62, w - 0.44, h - 0.8, body, size=bsize, color=MATN, line=1.16)


def stat(s, x, y, w, num, lab, col=ZARHAL):
    text(s, x, y, w, 0.86, num, size=40, color=col, font=H_FONT, bold=True,
         align=PP_ALIGN.CENTER)
    text(s, x, y + 0.86, w, 0.38, lab, size=11, color=SONIK, align=PP_ALIGN.CENTER)


def pic(s, name, x, y, w, folder="final", frame=True):
    p = shot(name, folder)
    if not p:
        box(s, x, y, w, w * 0.594, fill=PANEL, line=CHIZIQ)
        return
    if frame:
        box(s, x - 0.045, y - 0.045, w + 0.09, w * 0.594 + 0.09, fill=None, line=CHIZIQ, lw=1.0)
    s.shapes.add_picture(p, Inches(x), Inches(y), width=Inches(w))


def chevron(s, x, y, col=CHIZIQ, sz=0.13):
    """Kichik uchburchak — ketma-ketlik ko'rsatkichi (shrift belgisi emas)."""
    return box(s, x, y, sz, sz * 1.5, fill=col, shape=MSO_SHAPE.ISOSCELES_TRIANGLE)


def foot(s, n, txt=""):
    text(s, 0.72, 6.95, 8.0, 0.3, txt, size=9.5, color=CHIZIQ)
    text(s, 11.4, 6.95, 1.2, 0.3, "%02d" % n, size=9.5, color=CHIZIQ,
         align=PP_ALIGN.RIGHT)


# ===========================================================================
# 01 — Sarlavha
# ===========================================================================
s = new()
lg = os.path.join(BRAND, "logo_full.png")
s.shapes.add_picture(lg, Inches((W - 5.1) / 2), Inches(0.45), width=Inches(5.1))
text(s, 1.0, 5.35, 11.33, 0.45,
     "XIII asr Anadolusi  ·  Uchinchi shaxs  ·  Yashirinlik va jang",
     size=15, color=SONIK, align=PP_ALIGN.CENTER)
box(s, 5.4, 5.98, 2.53, 0.02, fill=CHIZIQ)
text(s, 1.0, 6.22, 11.33, 0.4,
     "48 epizod  ·  4 mavsum  ·  3 til  ·  o'z dvigateli, sof C++",
     size=13, color=ZARHAL, align=PP_ALIGN.CENTER)

# ===========================================================================
# 02 — Bir qarashda
# ===========================================================================
s = new(); y = title(s, "Bir qarashda", "Loyiha nima va nimasi bilan ajralib turadi")
text(s, 0.95, y, 7.1, 1.5,
     [[("Dirilish: Ertug'rul", {"bold": True, "color": ZARHAL, "size": 17}),
       (" — 1227-1261 yillar Anadolusida kechadigan uchinchi shaxs sarguzashti. "
        "Qayi obasining ko'chishi, Templar tartibi bilan to'qnashuv va Ertug'rulning "
        "shakllanishi 48 epizodda hikoya qilinadi.", {"size": 15})]],
     line=1.28)
text(s, 0.95, y + 1.35, 7.1, 1.2,
     "O'yin Assassin's Creed ning harakat erkinligini Sekiro ning jazolovchi jangi "
     "bilan birlashtiradi: yengil zarba yetarli emas, muvozanatni buzish kerak. "
     "Uchta dushman o'ldirishi mumkin, beshtasi — albatta.",
     size=14, color=SONIK, line=1.28)
for i, (n, l) in enumerate([("48", "epizod"), ("19.5k", "qator C++"),
                            ("921", "ovoz fayli"), ("928", "tarjima kaliti")]):
    stat(s, 0.95 + i * 1.82, y + 2.75, 1.72, n, l, ZARHAL if i % 2 == 0 else FERUZA)
pic(s, "valley_b.png", 8.55, y + 0.05, 3.85)
text(s, 8.55, y + 2.42, 3.85, 0.3, "Qayi obasi, vodiy — o'yin ichidan", size=10, color=CHIZIQ)
foot(s, 2, "BLADES OF ANATOLIA: ERTUGRUL")

# ===========================================================================
# 03 — Dunyo va davr
# ===========================================================================
s = new(); y = title(s, "Dunyo va davr", "Tarixiy langar: 1227-1261, Anadolu")
cards = [
    ("Qayi obasi", "To'rt yuz chodir, bir podadan boshqa boyligi yo'q ko'chmanchi oba. "
     "Ocharchilik va mo'g'ul bosimi ularni g'arbga siqib chiqaradi."),
    ("Uch kuch", "Salchuqiylar zaiflashmoqda, mo'g'ullar sharqdan bosmoqda, "
     "Templar tartibi esa savdo yo'llarini nazorat qilishga urinadi."),
    ("Ertug'rul", "Boshida — o'g'il. Oxirida — bey. O'yin uni qahramon deb "
     "ko'rsatmaydi: har qarori narx bilan keladi."),
]
for i, (h, b) in enumerate(cards):
    card(s, 0.95, y + i * 1.58, 5.5, 1.42, h, b, bsize=11.5)
pic(s, "09_cutscene.png", 6.85, y, 5.5)
box(s, 6.85, y + 3.46, 5.5, 1.42, fill=PANEL2, line=CHIZIQ, lw=0.75)
text(s, 7.08, y + 3.64, 5.05, 1.1,
     [[("Tarixiy aniqlik.", {"bold": True, "color": ZARHAL}),
       ("  Har epizodda hijriy va melodiy sana bor; hodisalar hujjatlashtirilganlik "
        "darajasi bilan belgilangan (DOCUMENTED / INFERRED / FICTION). Kodeks "
        "'buni qayerdan bilamiz' degan savolga javob beradi.", {})]],
     size=12, line=1.22)
foot(s, 3, "DUNYO")

# ===========================================================================
# 04 — Hikoya arxitekturasi
# ===========================================================================
s = new(); y = title(s, "Hikoya arxitekturasi", "4 mavsum x 12 epizod = 48 epizod, ~50 soat")
seasons = [
    ("S1", "1227-1229", "Och qarg'alar -> Bagras",
     "Ocharchilik, birinchi qon, oba ichidagi ishonchsizlik."),
    ("S2", "1230-1238", "Yassichemen -> Mix",
     "Katta janglar. Qo'l jarohati (Mix) — qaytarib bo'lmas narx."),
    ("S3", "1238-1243", "Qorda uch kun -> Sanoq",
     "Mo'g'ul bosqini. Omon qolish endi g'alabadan muhimroq."),
    ("S4", "1243-1261", "O'n ikki million -> Tanga",
     "Bey bo'lish. Qaror qilish — qilich urishdan og'irroq."),
]
yy = y
for i, (sid, yr, arc, note) in enumerate(seasons):
    box(s, 0.95, yy, 11.4, 1.02, fill=PANEL if i % 2 == 0 else PANEL2,
        line=CHIZIQ, lw=0.75)
    text(s, 1.22, yy + 0.20, 0.9, 0.5, sid, size=24, color=ZARHAL, font=H_FONT, bold=True)
    text(s, 2.20, yy + 0.16, 1.5, 0.3, yr, size=12, color=FERUZA, bold=True)
    text(s, 2.20, yy + 0.50, 3.4, 0.35, arc, size=12.5, color=SUYAK)
    text(s, 5.95, yy + 0.30, 6.15, 0.6, note, size=13, color=MATN, line=1.15)
    yy += 1.13
text(s, 0.95, yy + 0.12, 11.4, 0.4,
     [[("Har epizod", {"bold": True, "color": ZARHAL}),
       ("  ~62 daqiqa: intro cutscene -> brifing -> jang to'lqinlari -> "
        "cliffhanger. Keyingi epizod faqat oldingisi bajarilgach ochiladi.", {})]],
     size=13, color=MATN)
foot(s, 4, "HIKOYA")

# ===========================================================================
# 05 — Epizod arxetiplari
# ===========================================================================
s = new(); y = title(s, "Epizod arxetiplari",
                     "To'qqiz xil epizod turi — jang halqasi shundan quriladi")
arch = [
    ("SIEGE", "Qamal", "Ko'p to'lqin, tor joy"), ("DEFENSE", "Himoya", "Nuqtani ushlash"),
    ("SURVIVAL", "Omon qolish", "Vaqt va resurs"), ("INFILTRATION", "Kirib borish", "Sezilmaslik"),
    ("CHASE", "Quvish", "Nuqtaga yetish"), ("ESCORT", "Kuzatib borish", "Boshqani himoya qilish"),
    ("INVESTIGATION", "Tergov", "Izlash va so'roq"), ("COURT", "Kurultoy", "Muzokara va tanlov"),
    ("RITUAL", "Marosim", "Sekin, ramziy"),
]
for i, (en, uz, note) in enumerate(arch):
    cx = 0.95 + (i % 3) * 3.87
    cy = y + (i // 3) * 1.24
    box(s, cx, cy, 3.62, 1.10, fill=PANEL, line=CHIZIQ, lw=0.75)
    mark(s, cx + 0.24, cy + 0.30, 0.075, FERUZA if i % 2 == 0 else ZARHAL)
    text(s, cx + 0.44, cy + 0.20, 3.0, 0.3, en, size=11, color=FERUZA, bold=True)
    text(s, cx + 0.44, cy + 0.50, 3.0, 0.3, uz, size=14, color=SUYAK, font=H_FONT, bold=True)
    text(s, cx + 0.44, cy + 0.78, 3.0, 0.28, note, size=10.5, color=SONIK)
text(s, 0.95, y + 3.92, 11.4, 0.7,
     [[("Nima uchun muhim.", {"bold": True, "color": ZARHAL}),
       ("  episodes_v2.json da maqsadlar ro'yxati bo'sh edi. Endi maqsadlar "
        "ARXETIPDAN quriladi: DEFENSE -> barcha dushmanlarni yo'q qilish + omon qolish; "
        "INFILTRATION -> sezilmaslik (ixtiyoriy) + nuqtaga yetish + jang; "
        "CHASE -> nuqta + N dushman. To'lqinlar soni qiyinlikdan, dushmanlar soni "
        "max_simultaneous dan olinadi.", {})]],
     size=13, line=1.25)
foot(s, 5, "HIKOYA")

# ===========================================================================
# 06 — O'yin oqimi
# ===========================================================================
s = new(); y = title(s, "O'yin oqimi", "Ishga tushirishdan epizod yakuniga qadar")
steps = [
    ("01", "Splash", "Studiya belgisi"),
    ("02", "Til", "uz / tr / en"),
    ("03", "Bosh menyu", "Davom etish, epizodlar, sozlamalar"),
    ("04", "Epizod tanlash", "Ochilgani feruza, yopig'i so'nik"),
    ("05", "Cutscene", "Personajlar yuradi va gapiradi"),
    ("06", "Brifing", "Maqsadlar ko'rsatiladi"),
    ("07", "Jang", "To'lqinlar, nazorat nuqtalari"),
    ("08", "Natija", "Bajarildi yoki halok bo'ldingiz"),
]
for i, (n, h, b) in enumerate(steps):
    cx = 0.95 + (i % 4) * 2.90
    cy = y + (i // 4) * 1.42
    box(s, cx, cy, 2.66, 1.26, fill=PANEL, line=CHIZIQ, lw=0.75)
    text(s, cx + 0.22, cy + 0.16, 0.7, 0.3, n, size=13, color=FERUZA, font=H_FONT, bold=True)
    text(s, cx + 0.22, cy + 0.46, 2.3, 0.3, h, size=14, color=SUYAK, font=H_FONT, bold=True)
    text(s, cx + 0.22, cy + 0.76, 2.3, 0.3, b, size=10.5, color=SONIK, line=1.1)
    if i % 4 != 3:
        chevron(s, cx + 2.74, cy + 0.47, CHIZIQ, 0.12)
pic(s, "03_main.png", 0.95, y + 3.05, 3.62)
pic(s, "04_episodes.png", 4.84, y + 3.05, 3.62)
pic(s, "42_briefing.png", 8.73, y + 3.05, 3.62, folder="fin")
foot(s, 6, "OQIM")

# ===========================================================================
# 07 — Harakat mexanikasi
# ===========================================================================
s = new(); y = title(s, "Harakat: masofaga bog'langan qadam",
                     "Tizimning eng nozik qismi — va eng katta tuzatish")
text(s, 0.95, y, 6.5, 1.5,
     "Odatdagi protsedural yurishda qadam fazasi VAQTdan olinadi. Natijada tezlik "
     "o'zgarganda oyoq yerdan sirg'aladi. Bizda faza bosib o'tilgan MASOFAdan "
     "integrallanadi, tayanch fazasida panja esa dunyoda qotadi — oyoq burchaklari "
     "analitik IK bilan hisoblanadi.",
     size=14, color=MATN, line=1.28)
box(s, 0.95, y + 1.55, 6.5, 0.72, fill=PANEL, line=CHIZIQ, lw=0.75)
text(s, 1.18, y + 1.72, 6.1, 0.5,
     [[("S(v) = 0.97 + 0.35 x v", {"font": "Consolas", "color": ZARHAL, "size": 15, "bold": True})],
      [("faza += bosibOtilganMasofa / S(v)", {"font": "Consolas", "color": FERUZA, "size": 13})]],
     size=13, space=2)
rows = [("Tezlik", "Ilgari", "Endi", "Real odam", "SLIP"),
        ("Yurish 1.35 m/s", "4.6 qadam/s", "1.87", "1.85-1.90", "0.0%"),
        ("Yugurish 3.30 m/s", "4.6", "3.11", "3.0-3.2", "0.0%"),
        ("Sprint 6.40 m/s", "8.8", "3.99", "3.9-4.2", "0.0%")]
ty = y + 2.55
colx = [0.95, 3.30, 4.85, 6.05, 7.55]
colw = [2.30, 1.50, 1.15, 1.45, 1.10]
for ri, r in enumerate(rows):
    if ri == 0:
        box(s, 0.95, ty, 7.70, 0.36, fill=PANEL2, line=None)
    for ci, cell in enumerate(r):
        col = SONIK if ri == 0 else (MATN if ci == 0 else
              (YARA if ci == 1 else (ZARHAL if ci in (2, 4) else SUYAK)))
        text(s, colx[ci], ty + 0.09, colw[ci], 0.3, cell,
             size=10.5 if ri == 0 else 12,
             color=col, bold=(ri == 0 or ci in (2, 4)),
             align=PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.CENTER)
    ty += 0.42
text(s, 0.95, ty + 0.12, 7.7, 0.6,
     "Tezlik endi VEKTOR: yo'nalish bir kadrda o'zgarmaydi. Burilish tezligi "
     "fizikadan cheklanadi (a_lat <= 12 m/s^2) — sprintda 110 grad/s, yurishda 509.",
     size=12, color=SONIK, line=1.2)
pic(s, "63_sprint.png", 9.05, y + 0.05, 3.30, folder="fin")
pic(s, "69_arc2.png", 9.05, y + 2.10, 3.30, folder="loco")
text(s, 9.05, y + 4.05, 3.3, 0.3, "Sprint va burilish yoyi (tana egilgan)",
     size=10, color=CHIZIQ)
foot(s, 7, "MEXANIKA")

# ===========================================================================
# 08 — Parkur
# ===========================================================================
s = new(); y = title(s, "Parkur", "Assassin's Creed uslubidagi past va yuqori profil")
bullets(s, 0.95, y, 6.6, [
    ("Past profil.", "Yurish, cho'kkalash, olomonga singish. Shovqin deyarli yo'q."),
    ("Yuqori profil.", "Shift — yugurish va erkin yurish. To'siqlar avtomatik o'tiladi."),
    ("Zond fizikasi.", "probeVault, probeLedge, probeWall, mantleClear, supportBelow, "
     "platformAbove — oldindagi geometriya nur bilan tekshiriladi."),
    ("31 ta holat.", "Vault, Mantle, Climb, Hang, Shimmy, Eject, Slide, RollLand, "
     "WallRun, Dodge, LeapOfFaith va boshqalar."),
    ("Bilge Ko'z.", "H tugmasi — dunyo so'nadi, parkur geometriyasi yonadi."),
], size=13.5, gap=0.62)
pic(s, "60_idle.png", 7.85, y + 0.05, 4.50, folder="loco")
pic(s, "75_crouch.png", 7.85, y + 2.85, 4.50, folder="loco")
foot(s, 8, "MEXANIKA")

# ===========================================================================
# 09 — Jang: uch resurs
# ===========================================================================
s = new(); y = title(s, "Jang: uch resurs",
                     "GDD: jang og'ir va jazolovchi. Uch kishi o'ldirishi mumkin.")
res = [("Sog'liq", "Zarba zarari", "0 -> o'lim", YARA),
       ("Nafas", "Har harakat sarflaydi", "0 -> hujum yo'q, blok buziladi", SUYAK),
       ("Poza", "Zarba va blok to'playdi", "100 -> muvozanat buziladi", ZARHAL),
       ("Qo'l (Mix)", "Og'ir zarba va parry yeydi", "parry oynasi 180 -> 110 ms", FERUZA)]
for i, (n, d, f, c) in enumerate(res):
    cx = 0.95 + i * 2.90
    box(s, cx, y, 2.66, 1.55, fill=PANEL, line=CHIZIQ, lw=0.75)
    box(s, cx + 0.22, y + 0.30, 2.2, 0.10, fill=c)
    text(s, cx + 0.22, y + 0.52, 2.2, 0.3, n, size=15, color=c, font=H_FONT, bold=True)
    text(s, cx + 0.22, y + 0.84, 2.2, 0.3, d, size=11, color=MATN, line=1.1)
    text(s, cx + 0.22, y + 1.14, 2.2, 0.3, f, size=10, color=SONIK, line=1.1)
text(s, 0.95, y + 1.80, 6.6, 0.4, "Harakatlar", size=15, color=SUYAK,
     font=H_FONT, bold=True)
bullets(s, 0.95, y + 2.22, 6.6, [
    ("Yengil zarba.", "Uch zarbali kombo."),
    ("Og'ir zarba.", "Shift + Mouse 1 — o'z pozangizni ham buzadi."),
    ("Parry.", "Mouse 2 bosish, 180 ms oyna. Raqib pozasiga +35 — jangning yuragi."),
    ("Tepish.", "R — qalqonni ochadi. Blok qilib bo'lmaydi."),
    ("Yakunlovchi zarba.", "Muvozanati buzilgan raqibga — qat'iy o'lim."),
], size=13, gap=0.48)
pic(s, "46_fight4.png", 7.85, y + 1.85, 4.50, folder="fin")
foot(s, 9, "JANG")

# ===========================================================================
# 10 — Dushmanlar
# ===========================================================================
s = new(); y = title(s, "Dushmanlar", "Olti tur — har biri boshqacha yechim talab qiladi")
enemies = [
    ("Footman", "Oddiy piyoda", "Asosiy kuch. Yengil zarba bilan yechiladi."),
    ("Sergeant", "Qalqonli serjant", "Bloki kuchli — tepish kerak. O'q oldindan o'tmaydi."),
    ("Crossbow", "Arbaletchi", "Uzoqdan otadi, yaqinda zaif. Birinchi yo'q qiling."),
    ("Assassin", "Saroy qotili", "Tez, kam sog'liq. Parry qilish qiyin."),
    ("HorseArcher", "Mo'g'ul otliq", "S3+ da paydo bo'ladi. To'xtatib bo'lmaydi."),
    ("Elite", "Keshikten", "S4 elita gvardiyasi. Yakunlovchi zarba shart."),
]
for i, (n, r, d) in enumerate(enemies):
    cx = 0.95 + (i % 3) * 3.87
    cy = y + (i // 3) * 1.42
    box(s, cx, cy, 3.62, 1.28, fill=PANEL, line=CHIZIQ, lw=0.75)
    mark(s, cx + 0.24, cy + 0.28, 0.075, YARA)
    text(s, cx + 0.44, cy + 0.18, 3.0, 0.3, n, size=14, color=SUYAK, font=H_FONT, bold=True)
    text(s, cx + 0.44, cy + 0.50, 3.0, 0.28, r, size=10.5, color=FERUZA)
    text(s, cx + 0.24, cy + 0.80, 3.15, 0.4, d, size=11, color=MATN, line=1.12)
box(s, 0.95, y + 3.00, 11.4, 1.32, fill=PANEL2, line=CHIZIQ, lw=0.75)
text(s, 1.25, y + 3.22, 5.1, 0.9,
     [[("Sezish — ikki kanal.", {"bold": True, "color": ZARHAL}),
       ("  Ko'rish konusi (cho'kkalasangiz 55%, past profilda 85% masofa) va "
        "eshitish (parkur shovqinlaridan). Ogohlik bosh ustida to'lib boruvchi "
        "kvadrat bilan ko'rsatiladi.", {})]], size=12.5, line=1.2)
text(s, 6.75, y + 3.22, 5.3, 0.9,
     [[("Bir vaqtda faqat ikkitasi hujum qiladi.", {"bold": True, "color": ZARHAL}),
       ("  Qolganlari aylanib kutadi — AC dagi kabi. Aks holda besh kishilik "
        "to'lqin bir zumda o'ldirardi.", {})]], size=12.5, line=1.2)
foot(s, 10, "JANG")

# ===========================================================================
# 11 — Kamon
# ===========================================================================
s = new(); y = title(s, "Kamon", "XIII asr turk kompozit yoyi — ballistik o'q")
bullets(s, 0.95, y, 6.5, [
    ("Nishonga olish.", "G ni ushlang: kamera o'ng yelka ustiga siljiydi, "
     "ko'rish maydoni 48 dan 36 gradusga torayadi."),
    ("Tortish 0.85 s.", "Nafas sarflanadi (to'la tortib ushlaganda 11/s). Nafas "
     "kamaygach qo'l titraydi va nishon halqasi kengayadi."),
    ("Majburiy reliz.", "Nafas tugasa o'q o'zi uchadi — 60-70 kg yoyni uzoq ushlab bo'lmaydi."),
    ("Ballistika.", "O'q gravitatsiya bilan uchadi, 0.40 m qadamlar bilan "
     "integrallanadi; har qadamda yer, devor va dushman tekshiriladi."),
], size=13, gap=0.66)
zones = [("Oyoq", "x0.55", SONIK), ("Ko'krak", "x1.00", SUYAK), ("Bosh", "x2.60", YARA)]
text(s, 0.95, y + 2.85, 6.5, 0.35, "Tana zonalari", size=15, color=SUYAK,
     font=H_FONT, bold=True)
for i, (n, m, c) in enumerate(zones):
    cx = 0.95 + i * 2.20
    box(s, cx, y + 3.28, 1.96, 0.92, fill=PANEL, line=CHIZIQ, lw=0.75)
    text(s, cx, y + 3.42, 1.96, 0.4, m, size=22, color=c, font=H_FONT, bold=True,
         align=PP_ALIGN.CENTER)
    text(s, cx, y + 3.82, 1.96, 0.3, n, size=11, color=SONIK, align=PP_ALIGN.CENTER)
text(s, 0.95, y + 4.32, 6.5, 0.5,
     "Dubulg'asiz dushmanni boshdan bir o'q o'ldiradi — va u qichqirmaydi, "
     "ya'ni qo'shnilari bilmaydi. Devorga tekkan o'q qorovulni o'sha yoqqa yuboradi.",
     size=11.5, color=SONIK, line=1.2)
pic(s, "55_bow_full.png", 7.85, y + 0.05, 4.50, folder="fin")
pic(s, "56_bow_shot.png", 7.85, y + 2.85, 4.50, folder="fin")
foot(s, 11, "JANG")

# ===========================================================================
# 12 — Iymon
# ===========================================================================
s = new(); y = title(s, "Iymon", "Jangdagi joriy formangiz — yig'iladigan ochko emas")
tiers = [("Adashgan", "0-25", "nafas x0.75, poza x0.80", YARA),
         ("Shubha", "26-50", "neytral", SONIK),
         ("Sobit", "51-75", "nafas x1.50", ZARHAL),
         ("Sukunat", "76-100", "nafas/poza x1.50 + vaqt sekinlashuvi", FERUZA)]
for i, (n, r, e, c) in enumerate(tiers):
    cx = 0.95 + i * 2.90
    box(s, cx, y, 2.66, 1.62, fill=PANEL, line=c if i == 3 else CHIZIQ,
        lw=1.25 if i == 3 else 0.75)
    text(s, cx + 0.22, y + 0.20, 2.2, 0.3, r, size=11, color=SONIK)
    text(s, cx + 0.22, y + 0.52, 2.2, 0.35, n, size=17, color=c, font=H_FONT, bold=True)
    text(s, cx + 0.22, y + 0.94, 2.2, 0.55, e, size=11.5, color=MATN, line=1.15)
text(s, 0.95, y + 1.92, 5.5, 0.35, "Ko'tarish — mahoratdan", size=14, color=ZARHAL,
     font=H_FONT, bold=True)
bullets(s, 0.95, y + 2.32, 5.5, [
    "Mukammal parry  +0.4", "Yakunlovchi zarba  +2.0",
    "Zarba yemasdan tugatilgan to'lqin  +5.0", "Epizod bajarildi  +6.0",
], size=12.5, gap=0.36)
text(s, 6.85, y + 1.92, 5.5, 0.35, "Tushish — beparvolikdan", size=14, color=YARA,
     font=H_FONT, bold=True)
bullets(s, 6.85, y + 2.32, 5.5, [
    "Har zarba yeyish  -0.6", "Halok bo'lish  -8.0",
], size=12.5, gap=0.36)
box(s, 0.95, y + 3.90, 11.4, 0.98, fill=PANEL2, line=FERUZA, lw=1.0)
text(s, 1.25, y + 4.10, 10.8, 0.65,
     [[("Sukunat.", {"bold": True, "color": FERUZA, "size": 14}),
       ("  Iymon 76 dan yuqori bo'lsa, mukammal parry yoki yakunlovchi zarba "
        "1.2 soniya davomida vaqtni 0.45x ga sekinlashtiradi. 9 soniya sovish "
        "vaqti bor — uni yig'ib qo'yib bo'lmaydi. Iymon progress.json da saqlanadi.",
        {"size": 13})]], line=1.22)
foot(s, 12, "JANG")

# ===========================================================================
# 13 — Zarba qaytarmasi
# ===========================================================================
s = new(); y = title(s, "Zarba qaytarmasi",
                     "Nima uchun jang 'bo'sh' his qilinardi — uchta uzilgan sim")
box(s, 0.95, y, 11.4, 1.42, fill=PANEL, line=YARA, lw=1.0)
text(s, 1.25, y + 0.20, 10.8, 1.05,
     [[("Topilgan asosiy nuqson.", {"bold": True, "color": YARA, "size": 14}),
       ("  setClip() bir martalik klipni HAR KADR qayta boshlardi. Character va Enemy "
        "klipni har kadr so'raydi, setClip esa mahalliy vaqtni nolga qaytarardi — "
        "natijada zarba, zarba yeyish, parkur va yakunlovchi zarba animatsiyalari "
        "birinchi kadrida MUZLAB turardi. Poza lt = dt da qotgani uchun natija kadr "
        "chastotasiga ham bog'liq edi: 60 Gts da bosh 21 gradus, 30 Gts da 33 gradus "
        "orqaga tashlangan holda 0.35 s qotardi.", {"size": 12.5})]], line=1.22)
hdr = ["Kanal", "Tegdi", "Bloklandi", "O'ldirdi", "O'zi yedi", "Parry"]
data = [["Muzlash", "55 ms", "35 ms", "100 ms", "70 ms", "130 ms"],
        ["Kamera turtkisi", "0.28", "0.18", "0.45", "0.65", "0.40"],
        ["Uchqun", "feruza", "zarhal", "ikki qavat", "feruza", "zarhal"],
        ["Tovush", "past chuq", "metall", "og'ir", "kuchli", "yorqin"]]
cx0 = [0.95, 3.35, 4.95, 6.60, 8.35, 10.15]
cw = [2.35, 1.55, 1.60, 1.70, 1.75, 1.60]
ty = y + 1.72
box(s, 0.95, ty, 11.4, 0.38, fill=PANEL2, line=None)
for ci, hh in enumerate(hdr):
    text(s, cx0[ci], ty + 0.10, cw[ci], 0.3, hh, size=11, color=SONIK, bold=True,
         align=PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.CENTER)
ty += 0.44
for row in data:
    for ci, cell in enumerate(row):
        text(s, cx0[ci], ty, cw[ci], 0.3, cell, size=12,
             color=MATN if ci == 0 else (ZARHAL if ci == 4 else SUYAK),
             bold=(ci == 4),
             align=PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.CENTER)
    ty += 0.40
text(s, 0.95, ty + 0.12, 11.4, 0.8,
     [[("Ikki nozik joy.", {"bold": True, "color": ZARHAL}),
       ("  Muzlash HAQIQIY dt bilan so'nadi — aks holda o'zi sekinlashtirgan vaqt "
        "bilan kamayib cheksiz cho'zilardi. Kamera silkinishi camSmooth ga "
        "yozilmaydi — yozilsa turtkilar qo'shilib kamerani joyidan surib yuborardi.",
        {})]], size=12.5, line=1.22)
foot(s, 13, "JANG")

# ===========================================================================
# 14 — Epizod halqasi
# ===========================================================================
s = new(); y = title(s, "Epizod halqasi", "Muvaffaqiyatsizlik va qayta tug'ilishning to'liq mantig'i")
flow = [("Cutscene", FERUZA), ("Brifing", FERUZA), ("Jang", ZARHAL),
        ("To'lqin tozalandi", ZARHAL), ("Nazorat nuqtasi", FERUZA)]
for i, (t, c) in enumerate(flow):
    cx = 0.95 + i * 2.42
    box(s, cx, y, 2.18, 0.62, fill=PANEL, line=c, lw=0.9)
    text(s, cx, y + 0.19, 2.18, 0.3, t, size=12, color=SUYAK, bold=True,
         align=PP_ALIGN.CENTER)
    if i < 4:
        chevron(s, cx + 2.24, y + 0.22, CHIZIQ, 0.12)
box(s, 0.95, y + 0.85, 5.5, 1.72, fill=PANEL, line=YARA, lw=1.0)
text(s, 1.20, y + 1.02, 5.0, 0.35, "Halok bo'ldingiz", size=16, color=YARA,
     font=H_FONT, bold=True)
bullets(s, 1.20, y + 1.42, 5.0, [
    "Nazorat nuqtasidan qayta urinish",
    "Epizodni boshidan",
    "Bosh menyuga",
], size=12, gap=0.34)
text(s, 1.20, y + 2.52, 5.0, 0.3,
     "Joriy to'lqin dushmanlari yangidan, bajarilganlari qaytmaydi.",
     size=10.5, color=SONIK)
box(s, 6.85, y + 0.85, 5.5, 1.72, fill=PANEL, line=FERUZA, lw=1.0)
text(s, 7.10, y + 1.02, 5.0, 0.35, "Epizod bajarildi", size=16, color=FERUZA,
     font=H_FONT, bold=True)
bullets(s, 7.10, y + 1.42, 5.0, [
    "Cliffhanger cutscene",
    "Keyingi epizod ochiladi",
    "saves/progress.json ga yoziladi",
], size=12, gap=0.34)
text(s, 7.10, y + 2.52, 5.0, 0.3,
     "O'limlar, o'ldirilganlar va Iymon saqlanadi.",
     size=10.5, color=SONIK)
pic(s, "49_result.png", 2.15, y + 2.72, 4.0, folder="fin")
pic(s, "51_respawn.png", 7.15, y + 2.72, 4.0, folder="fin")
foot(s, 14, "OQIM")

# ===========================================================================
# 15 — Interfeys
# ===========================================================================
s = new(); y = title(s, "Interfeys: Temir va Firuza",
                     "O'ziga xos dizayn tizimi — chapga tekislangan, muharrirona")
pal = [("#0E1316", "fon", FON), ("#161D21", "panel", PANEL), ("#2A353A", "chiziq", CHIZIQ),
       ("#E4EAEA", "matn", MATN), ("#7C8B8F", "so'nik", SONIK), ("#48A9B5", "feruza", FERUZA),
       ("#C09660", "zarhal", ZARHAL), ("#BC5A44", "yara", YARA), ("#DCD3C4", "suyak", SUYAK)]
for i, (hexv, nm, c) in enumerate(pal):
    cx = 0.95 + i * 1.28
    box(s, cx, y, 1.12, 0.72, fill=c, line=CHIZIQ, lw=0.75)
    text(s, cx, y + 0.80, 1.12, 0.25, nm, size=10, color=SUYAK, align=PP_ALIGN.CENTER)
    text(s, cx, y + 1.04, 1.12, 0.25, hexv, size=8.5, color=CHIZIQ, align=PP_ALIGN.CENTER)
bullets(s, 0.95, y + 1.55, 6.2, [
    ("Kichik kvadrat.", "Butun interfeysning yagona motivi — 'mix boshi'."),
    ("Feruza = tanlash.", "Zarhal = qiymat. Yara = xavf. Boshqa rang yo'q."),
    ("Klavish qayta biriktirish.", "Har harakat ustiga bosing va yangi klavishni "
     "bosing — sozlamalar bosilgan klavish NOMINI ko'rsatadi."),
    ("Uch til.", "Menyu, HUD, subtitr va epizod nomlari — hammasi tarjima qilingan."),
], size=13, gap=0.60)
pic(s, "05_settings.png", 7.50, y + 1.50, 4.85)
foot(s, 15, "INTERFEYS")

# ===========================================================================
# 16 — Audio
# ===========================================================================
s = new(); y = title(s, "Audio", "912 ta ovozli replika + 9 ta protsedural jang tovushi")
card(s, 0.95, y, 5.5, 1.80, "Ovozli replikalar",
     "Har uch til uchun 309 tadan WAV — jami 912. PowerShell System.Speech bilan "
     "oldindan yaratilgan, o'yin ichida SAPI zaxira yo'li ham bor. "
     "Cutscene subtitri bilan aniq sinxron.", ZARHAL)
card(s, 6.85, y, 5.5, 1.80, "Aralashtirgich",
     "waveOut asosidagi o'z mikseri: 5 shina (master, musiqa, effekt, ovoz, muhit). "
     "Tashqi kutubxona yo'q. Sozlamalarda har shina alohida boshqariladi.", FERUZA)
text(s, 0.95, y + 2.02, 11.4, 0.4,
     "Jang tovushlari — protsedural sintez, tovush fayli yo'q",
     size=16, color=SUYAK, font=H_FONT, bold=True)
sfx = [("O'lim", "past pasayuvchi sinus + shovqin", "1108 Gts"),
       ("Zarba", "190->72 Gts + quruq shovqin", "1922 Gts"),
       ("Qalqon", "garmonik BO'LMAGAN qismlar 1 : 1.52 : 2.14 : 2.84", "2227 Gts"),
       ("Parry", "o'sha nisbat, yuqoriroq va uzunroq ring", "3291 Gts"),
       ("Qilich yoyi", "o'rtada ochiladigan filtr + qo'ng'iroq konvert", "3524 Gts")]
ty = y + 2.48
for i, (n, d, f) in enumerate(sfx):
    box(s, 0.95, ty, 11.4, 0.46, fill=PANEL if i % 2 == 0 else PANEL2, line=None)
    mark(s, 1.20, ty + 0.18, 0.075, FERUZA)
    text(s, 1.42, ty + 0.11, 2.0, 0.28, n, size=12, color=SUYAK, bold=True)
    text(s, 3.55, ty + 0.11, 6.7, 0.28, d, size=11.5, color=MATN)
    text(s, 10.40, ty + 0.11, 1.85, 0.28, f, size=11.5, color=ZARHAL,
         align=PP_ALIGN.RIGHT, bold=True)
    ty += 0.50
text(s, 0.95, ty + 0.10, 11.4, 0.35,
     "Metall jarangi BUTUN SONLI BO'LMAGAN nisbatdan chiqadi — garmonik qismlar "
     "'musiqiy' eshitiladi va zarbaga o'xshamaydi.",
     size=11.5, color=SONIK, line=1.18)
foot(s, 16, "AUDIO")

# ===========================================================================
# 17 — Texnologiya
# ===========================================================================
s = new(); y = title(s, "Texnologiya", "O'z dvigateli — tashqi kutubxona nol ta")
left = [("Til", "C++20, MinGW-w64 g++ 16.1"),
        ("Grafika", "Win32 + fixed-function OpenGL 1.1 + GLU"),
        ("Yig'ish", "CMake + Ninja, 25 nishon"),
        ("Rasm", "GDI+ dekod -> GL teksturalari"),
        ("Shrift", "GDI glif atlasi (UTF-8, uz/tr belgilari)"),
        ("Ovoz", "waveOut mikser + SAPI")]
right = [("GLFW", "yo'q"), ("GLM", "yo'q"), ("tinyobjloader", "yo'q"),
         ("Shader / VBO", "yo'q"), ("Fizika dvigateli", "yo'q"),
         ("nlohmann/json", "yagona bundle qilingan header")]
text(s, 0.95, y, 5.6, 0.35, "Nima ishlatiladi", size=15, color=ZARHAL,
     font=H_FONT, bold=True)
ty = y + 0.45
for k, v in left:
    box(s, 0.95, ty, 5.5, 0.50, fill=PANEL, line=None)
    text(s, 1.18, ty + 0.13, 1.9, 0.3, k, size=12, color=FERUZA, bold=True)
    text(s, 3.10, ty + 0.13, 3.2, 0.3, v, size=12, color=MATN)
    ty += 0.55
text(s, 6.85, y, 5.5, 0.35, "Nima ishlatilmaydi", size=15, color=SONIK,
     font=H_FONT, bold=True)
ty = y + 0.45
for k, v in right:
    box(s, 6.85, ty, 5.5, 0.50, fill=PANEL, line=None)
    text(s, 7.08, ty + 0.13, 2.6, 0.3, k, size=12, color=SUYAK, bold=True)
    text(s, 9.75, ty + 0.13, 2.5, 0.3, v, size=12,
         color=YARA if v == "yo'q" else SONIK)
    ty += 0.55
box(s, 0.95, y + 3.90, 11.4, 1.05, fill=PANEL2, line=CHIZIQ, lw=0.75)
text(s, 1.25, y + 4.10, 10.8, 0.7,
     [[("Modellar statik .obj — skelet yo'q.", {"bold": True, "color": ZARHAL}),
       ("  Skin.cpp chegara qutisi bo'yicha 12 ta virtual suyakka avtomatik rigging "
        "qiladi va 30 Gts da CPU da skinlaydi. Bo'yin balandligi meshdan o'lchanadi, "
        "36 ta animatsiya protsedural poza funksiyalari bilan quriladi.", {})]],
     size=12.5, line=1.22)
foot(s, 17, "TEXNOLOGIYA")

# ===========================================================================
# 18 — Raqamlar
# ===========================================================================
s = new(); y = title(s, "Raqamlar", "Loyihaning hozirgi holati")
nums = [("19 528", "qator C++"), ("24", "manba fayli"), ("48", "epizod"),
        ("48", "cutscene"), ("4", "daraja"), ("921", "audio fayl"),
        ("928", "tarjima kaliti"), ("36", "animatsiya klipi"),
        ("31", "harakat holati"), ("6", "dushman turi"),
        ("9", "epizod arxetipi"), ("3", "til")]
for i, (n, l) in enumerate(nums):
    cx = 0.95 + (i % 6) * 1.92
    cy = y + (i // 6) * 1.65
    box(s, cx, cy, 1.74, 1.42, fill=PANEL, line=CHIZIQ, lw=0.75)
    text(s, cx, cy + 0.28, 1.74, 0.6, n, size=26, color=ZARHAL if i % 2 == 0 else FERUZA,
         font=H_FONT, bold=True, align=PP_ALIGN.CENTER)
    text(s, cx, cy + 0.92, 1.74, 0.35, l, size=10.5, color=SONIK, align=PP_ALIGN.CENTER)
box(s, 0.95, y + 3.55, 11.4, 0.95, fill=PANEL2, line=FERUZA, lw=1.0)
text(s, 1.25, y + 3.74, 10.8, 0.6,
     [[("Sifat holati.", {"bold": True, "color": FERUZA}),
       ("  Toza yig'ilish: 25/25 nishon, 0 xato, 0 ogohlantirish. "
        "Kontent tekshiruvi: 928 kalitdan 0 tasi topilmagan, 921 audio fayl joyida. "
        "Oyoq sirg'alishi: SLIP 0.0% (mezon 3%).", {})]], size=13, line=1.22)
foot(s, 20, "HOLAT")

# ===========================================================================
# 19 — Ochiq kamchiliklar
# ===========================================================================
s = new(); y = title(s, "Ochiq kamchiliklar", "Halol ro'yxat — hali bajarilmagan ishlar")
gaps = [
    ("Jang balansi", "Vitals::receive poza zararini 0.7 ga qisadi; zarba jadvali "
     "shunga qarab qayta balanslanmagan.", YARA),
    ("Parry bepul", "Blok nafas sarflaydi, parry esa yo'q — muvozanatsizlik.", YARA),
    ("Dushman bloki", "Burchak tekshirilmaydi: orqadan urganda ham qalqon ushlaydi.", YARA),
    ("Zarba balandligi", "O'yinchi zarbasi Y o'qini tekshirmaydi — pastdagi va "
     "tepadagi dushmanga ham tegadi.", YARA),
    ("Skelet animatsiya", "Modellarda suyak yo'q; quti bo'yicha rigging chegarasi bor.", SONIK),
    ("Shader yo'q", "Fixed-function OpenGL 1.1 — soya xaritasi va PBR imkonsiz.", SONIK),
    ("Faqat Windows", "Win32 ga bevosita bog'langan; port qilish katta ish.", SONIK),
    ("Git Bash segfault", "PowerShell va Explorer dan ishlaydi, Git Bash dan yo'q — "
     "sabab aniqlanmagan.", SONIK),
]
for i, (h, b, c) in enumerate(gaps):
    cx = 0.95 + (i % 2) * 5.90
    cy = y + (i // 2) * 1.22
    box(s, cx, cy, 5.55, 1.08, fill=PANEL, line=CHIZIQ, lw=0.75)
    mark(s, cx + 0.22, cy + 0.26, 0.075, c)
    text(s, cx + 0.42, cy + 0.16, 4.9, 0.3, h, size=13, color=SUYAK,
         font=H_FONT, bold=True)
    text(s, cx + 0.22, cy + 0.50, 5.1, 0.5, b, size=11, color=MATN, line=1.15)
foot(s, 21, "HOLAT")

# ===========================================================================
# 18 — Demo xarita va aktivlar
# ===========================================================================
s = new(); y = title(s, "Demo xarita va aktivlar",
                     "Xaritaning o'rtasi Blender'da tayyorlanmoqda")
box(s, 0.95, y, 5.6, 1.95, fill=PANEL, line=FERUZA, lw=1.0)
mark(s, 1.20, y + 0.28, 0.085, FERUZA)
text(s, 1.42, y + 0.18, 4.9, 0.35, "Demo xarita", size=17, color=FERUZA,
     font=H_FONT, bold=True)
text(s, 1.20, y + 0.66, 5.1, 1.5,
     "Yangi demo xarita — \"Qayi obasi, vodiy\": halqa bo'lib tizilgan o'tovlar, "
     "tepaliklar, qalin qarag'ayzor va oltin soat yoritishi. 1 884 obyekt, "
     "1 457 to'qnashuv qutisi. Markaziy qism Blender'da modellashtirilmoqda.",
     size=13, color=MATN, line=1.22)

box(s, 6.75, y, 5.6, 1.95, fill=PANEL, line=ZARHAL, lw=1.0)
mark(s, 7.00, y + 0.28, 0.085, ZARHAL)
text(s, 7.22, y + 0.18, 4.9, 0.35, "Chodir (o'tov) — tayyor", size=17, color=ZARHAL,
     font=H_FONT, bold=True)
text(s, 7.00, y + 0.66, 5.1, 1.5,
     "Qayi o'tovi modellashtirildi va teksturalandi: 48 565 uchburchak, "
     "27 280 verteks, uch burchakli topologiya. Model tayyor va o'yinga "
     "qo'yilish arafasida turibdi — keyingi qadam shu.",
     size=13, color=MATN, line=1.22)

# O'tov renderi uchun joy (docs/brand/yurt.png bo'lsa — qo'yiladi)
_yurt = os.path.join(BRAND, "yurt.png")
if os.path.exists(_yurt):
    box(s, 0.90, y + 2.18, 5.70, 2.62, fill=None, line=CHIZIQ, lw=1.0)
    s.shapes.add_picture(_yurt, Inches(0.95), Inches(y + 2.23), height=Inches(2.52))
else:
    box(s, 0.95, y + 2.18, 5.6, 2.55, fill=PANEL2, line=CHIZIQ, lw=0.75)
    text(s, 1.20, y + 3.10, 5.1, 0.6,
         [[("O'tov renderi shu yerga qo'yiladi", {"bold": True, "color": SONIK}),
           ("\ndocs/brand/yurt.png", {"color": CHIZIQ, "font": "Consolas", "size": 11})]],
         size=13, color=SONIK, align=PP_ALIGN.CENTER)
pic(s, "valley_a.png", 6.80, y + 2.23, 4.30)
text(s, 6.80, y + 4.85, 5.5, 0.3, "Qayi obasi, vodiy — o'yin ichidan",
     size=10, color=CHIZIQ)
foot(s, 20, "ISHLAB CHIQISH")

# ===========================================================================
# 19 — Ishlab chiqish holati
# ===========================================================================
s = new(); y = title(s, "Ishlab chiqish holati",
                     "Loyiha boshlanganiga bir hafta bo'ldi")
for i, (nn, ll) in enumerate([("1", "hafta"), ("48", "epizod"),
                              ("19.5k", "qator C++"), ("1", "demo xarita")]):
    box(s, 0.95 + i * 2.90, y, 2.66, 1.45, fill=PANEL, line=CHIZIQ, lw=0.75)
    text(s, 0.95 + i * 2.90, y + 0.22, 2.66, 0.82, nn, size=36,
         color=ZARHAL if i % 2 == 0 else FERUZA, font=H_FONT, bold=True,
         align=PP_ALIGN.CENTER)
    text(s, 0.95 + i * 2.90, y + 0.96, 2.66, 0.35, ll, size=11.5, color=SONIK,
         align=PP_ALIGN.CENTER)

box(s, 0.95, y + 1.75, 11.4, 1.42, fill=PANEL, line=ZARHAL, lw=1.0)
mark(s, 1.22, y + 2.02, 0.085, ZARHAL)
text(s, 1.45, y + 1.92, 10.6, 1.15,
     [[("Aktivlar sun'iy intellekt bilan yaratilgan.", {"bold": True, "color": ZARHAL, "size": 14}),
       ("  Modellar, teksturalar va ovozli replikalarning katta qismi AI orqali "
        "generatsiya qilindi. Aynan shuning uchun ularni o'yinga moslashtirish "
        "vaqt oladi: har bir model qayta topologiyalanishi, masshtabi va o'qlari "
        "to'g'rilanishi, teksturasi o'yin uslubiga keltirilishi kerak. "
        "Chodir ana shu jarayondan o'tdi va endi o'yinga qo'yiladi.", {"size": 13})]],
     line=1.24)

text(s, 0.95, y + 3.40, 11.4, 0.4, "Bir haftada nimalar qilindi",
     size=16, color=SUYAK, font=H_FONT, bold=True)
done = [
    ("Yadro", "O'z dvigateli, menyular, 3 til, 48 epizod ma'lumoti"),
    ("Harakat", "AC darajasidagi lokomotsiya, parkur, oyoq sirg'alishi 0%"),
    ("Jang", "Uch resurs, 6 dushman, kamon, Iymon, zarba qaytarmasi"),
    ("Halqa", "To'lqinlar, nazorat nuqtasi, o'lim va qayta tug'ilish"),
]
for i, (h, b) in enumerate(done):
    cx = 0.95 + (i % 2) * 5.90
    cy = y + 3.85 + (i // 2) * 0.92
    box(s, cx, cy, 5.55, 0.80, fill=PANEL, line=None)
    mark(s, cx + 0.22, cy + 0.20, 0.075, FERUZA)
    text(s, cx + 0.44, cy + 0.10, 1.5, 0.3, h, size=12.5, color=FERUZA, bold=True)
    text(s, cx + 1.95, cy + 0.10, 3.4, 0.6, b, size=11.5, color=MATN, line=1.15)
foot(s, 21, "ISHLAB CHIQISH")

# ===========================================================================
# 22 — Loyiha videosi
# ===========================================================================
# 22 — Loyiha videosi

# ===========================================================================
s = new(); y = title(s, "Loyiha videosi",
                     "2 daqiqa 55 soniya — o'yinning hozirgi yakuniy ko'rinishi")
pic(s, "video_still.png", 0.95, y, 7.3)
box(s, 8.55, y, 3.80, 4.34, fill=PANEL, line=FERUZA, lw=1.0)
text(s, 8.80, y + 0.22, 3.3, 0.35, "Videoda nima bor", size=15, color=FERUZA,
     font=H_FONT, bold=True)
bullets(s, 8.80, y + 0.70, 3.3, [
    "Vodiy xaritasi — kran plani",
    "Yurish, yugurish, sprint",
    "Burilish yoyi va inersiya",
    "Sakrash va parkur",
    "Bilge Ko'z",
    "Kamon — tortish va otish",
    "Qilich jangi va parry",
    "Halokat va qayta tug'ilish",
], size=11.5, gap=0.42)
box(s, 0.95, y + 4.62, 11.4, 0.80, fill=PANEL2, line=CHIZIQ, lw=0.75)
text(s, 1.22, y + 4.80, 10.9, 0.5,
     [[("Havola:", {"bold": True, "color": ZARHAL}),
       ("   docs/Blades_of_Anatolia_demo.mp4   —   1280x720, 30 k/s, 33 MB",
        {"font": "Consolas", "size": 12.5, "color": MATN})]], size=13)
foot(s, 22, "VIDEO")

# ===========================================================================
# 23 — Jamoa
# ===========================================================================
s = new(); y = title(s, "Jamoa", "Uch kishi")
team = [
    ("MANSUROV ISLOMBEK BAXROMBEKOVICH", "Jamoa rahbari",
     "Game Developer / Designer", "+998 99 790 15 53",
     "mansurovislombek130@gmail.com", "167-sonli maktab", ZARHAL),
    ("RAXMANOVA MARFU'A YUSUFBEKOVNA", "A'zo",
     "Co-Founder", "+998 99 790 15 53",
     "rakhmonovamunira85@gmail.com",
     "Respublika Ixtisoslashtirilgan Ftiziatriya va Pulmonologiya Markazi", FERUZA),
    ("AGZAMOV MIRAKBAR ALISHER O'G'LI", "A'zo",
     "AI muhandis", "+998 99 202 57 97",
     "agzamovanozimahon@gmail.com", "280-sonli maktab", SUYAK),
]
yy = y
for (nm, rol, kasb, tel, mail, joy, col) in team:
    box(s, 0.95, yy, 11.4, 1.48, fill=PANEL, line=CHIZIQ, lw=0.75)
    box(s, 1.22, yy + 0.30, 0.86, 0.86, fill=PANEL2, line=col, lw=1.0)
    text(s, 1.22, yy + 0.52, 0.86, 0.52, nm[0], size=24, color=col,
         font=H_FONT, bold=True, align=PP_ALIGN.CENTER)
    text(s, 2.32, yy + 0.22, 6.4, 0.34, nm, size=14, color=SUYAK,
         font=H_FONT, bold=True)
    text(s, 2.32, yy + 0.60, 6.4, 0.30, rol + "  ·  " + kasb, size=12, color=col)
    text(s, 2.32, yy + 0.94, 6.4, 0.30, joy, size=11, color=SONIK)
    text(s, 8.95, yy + 0.34, 3.2, 0.30, tel, size=11.5, color=MATN,
         align=PP_ALIGN.RIGHT)
    text(s, 8.45, yy + 0.68, 3.7, 0.30, mail, size=11.5, color=SONIK,
         align=PP_ALIGN.RIGHT)
    yy += 1.60
foot(s, 23, "JAMOA")

# ===========================================================================
# 24 — Yakun
# ===========================================================================
s = new()
s.shapes.add_picture(os.path.join(BRAND, "logo_mark.png"),
                     Inches((W - 2.4) / 2), Inches(1.05), width=Inches(2.4))
text(s, 1.0, 3.72, 11.33, 0.6, "Bir jumlada", size=15, color=FERUZA,
     align=PP_ALIGN.CENTER, font=H_FONT, bold=True)
text(s, 2.0, 4.20, 9.33, 1.2,
     "Tashqi kutubxonasiz, sof C++ da yozilgan uchinchi shaxs o'yini: "
     "AC darajasidagi harakat, Sekiro og'irligidagi jang va 48 epizodli "
     "tarixiy hikoya — hammasi bitta dvigatelda.",
     size=17, color=MATN, align=PP_ALIGN.CENTER, line=1.35, font=H_FONT)
box(s, 5.4, 5.72, 2.53, 0.02, fill=CHIZIQ)
text(s, 1.0, 5.98, 11.33, 0.4,
     "19 528 qator  ·  25/25 nishon  ·  0 xato  ·  0 ogohlantirish",
     size=13, color=ZARHAL, align=PP_ALIGN.CENTER)

out = os.path.join(HERE, "Blades_of_Anatolia_Ertugrul.pptx")
prs.save(out)
print("saqlandi:", out, "slaydlar:", len(prs.slides.__iter__.__self__._sldIdLst))
